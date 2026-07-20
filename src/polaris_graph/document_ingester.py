"""
Local document ingestion for POLARIS sovereign/air-gap deployments.

Parses uploaded files LOCALLY -- never sends content to external APIs.
Supports: PDF, DOCX, DOC, XLSX, PPTX, TXT, MD, CSV, HTML, HTM,
          Images (PNG, JPG, JPEG, WEBP, GIF), Audio (MP3, WAV, M4A, AAC, OGG, FLAC).

All extraction is performed by local libraries:
  - PDF:  PyMuPDF (fitz), with pytesseract OCR fallback for vector-curve PDFs
  - DOCX: python-docx
  - DOC:  antiword (preferred) or raw text extraction fallback
  - XLSX: openpyxl
  - PPTX: python-pptx
  - HTML/HTM: readability-lxml
  - CSV:  stdlib csv module
  - TXT/MD: direct read
  - Images: pytesseract OCR + PIL for metadata
  - Audio: OpenAI Whisper for transcription

Storage layout per document:
  data/documents/{doc_id}/
    original.*            # Original uploaded file
    extracted.txt         # Plain-text extraction
    extracted.html        # HTML-formatted extraction
    metadata.json         # Extraction metadata
"""

import csv
import hashlib
import io
import json
import logging
import os
import shutil
import subprocess
import time
from pathlib import Path
from typing import Optional

from dotenv import load_dotenv
from src.polaris_graph.settings import resolve

load_dotenv()

logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Configuration (LAW VI: all from env vars)
# ---------------------------------------------------------------------------

MAX_UPLOAD_SIZE_MB = int(resolve("PG_MAX_UPLOAD_SIZE_MB"))
DOCUMENT_STORAGE_DIR = Path(resolve("PG_DOCUMENT_STORAGE_DIR"))

# OCR fallback configuration for vector-curve PDFs (CorelDRAW, Illustrator, etc.)
TESSERACT_CMD = os.getenv(
    "TESSERACT_CMD",
    r"C:\Program Files\Tesseract-OCR\tesseract.exe",
)
OCR_DPI = int(resolve("PG_OCR_DPI"))
OCR_MAX_PAGES = int(resolve("PG_OCR_MAX_PAGES"))
OCR_MIN_TEXT_CHARS = int(resolve("PG_OCR_MIN_TEXT_CHARS"))

# Audio transcription configuration
WHISPER_MODEL = os.getenv("WHISPER_MODEL", "base")
FFMPEG_PATH = os.getenv("FFMPEG_PATH", r"C:\Windows\ffmpeg.exe")

# Legacy .doc parsing: path to antiword binary (if available)
ANTIWORD_CMD = os.getenv("ANTIWORD_CMD", "antiword")

# Image extensions (OCR via pytesseract)
_IMAGE_EXTENSIONS = frozenset({".png", ".jpg", ".jpeg", ".webp", ".gif"})

# Audio extensions (transcription via Whisper)
_AUDIO_EXTENSIONS = frozenset({".mp3", ".wav", ".m4a", ".aac", ".ogg", ".flac"})

ACCEPTED_EXTENSIONS = frozenset({
    ".pdf", ".docx", ".doc", ".xlsx", ".pptx",
    ".txt", ".md", ".csv", ".html", ".htm",
}) | _IMAGE_EXTENSIONS | _AUDIO_EXTENSIONS


class DocumentIngestionError(Exception):
    """Raised when document ingestion fails for a recoverable reason."""


class DocumentIngester:
    """Parse uploaded files locally and persist extracted content.

    Every parse method returns a dict with keys:
        content  (str):  Plain-text extraction.
        html     (str):  HTML-formatted extraction.
        metadata (dict): Format-specific metadata (author, page count, etc.).
        pages    (int):  Number of pages/sheets/slides (1 for flat text).
    """

    # Extension -> method name mapping.  Methods are resolved via getattr
    # so they must exist on the class as ``_parse_{ext}`` (without the dot).
    PARSERS: dict[str, str] = {
        ".pdf":  "_parse_pdf",
        ".docx": "_parse_docx",
        ".doc":  "_parse_doc",
        ".xlsx": "_parse_xlsx",
        ".pptx": "_parse_pptx",
        ".txt":  "_parse_text",
        ".md":   "_parse_text",
        ".csv":  "_parse_csv",
        ".html": "_parse_html",
        ".htm":  "_parse_html",
        # Image formats (OCR)
        ".png":  "_parse_image",
        ".jpg":  "_parse_image",
        ".jpeg": "_parse_image",
        ".webp": "_parse_image",
        ".gif":  "_parse_image",
        # Audio formats (Whisper transcription)
        ".mp3":  "_parse_audio",
        ".wav":  "_parse_audio",
        ".m4a":  "_parse_audio",
        ".aac":  "_parse_audio",
        ".ogg":  "_parse_audio",
        ".flac": "_parse_audio",
    }

    # ------------------------------------------------------------------
    # Public API
    # ------------------------------------------------------------------

    async def ingest(self, file_path: Path) -> dict:
        """Ingest a single file and persist extracted artefacts.

        Args:
            file_path: Absolute or relative path to the uploaded file.

        Returns:
            dict with keys: content, html, metadata, pages, doc_id.

        Raises:
            DocumentIngestionError: On validation or parse failures.
            FileNotFoundError:      If *file_path* does not exist.
        """
        file_path = Path(file_path).resolve()

        # --- validation ---------------------------------------------------
        if not file_path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")

        ext = file_path.suffix.lower()
        if ext not in ACCEPTED_EXTENSIONS:
            raise DocumentIngestionError(
                f"Unsupported file extension '{ext}'. "
                f"Accepted: {sorted(ACCEPTED_EXTENSIONS)}"
            )

        size_mb = file_path.stat().st_size / (1024 * 1024)
        if size_mb > MAX_UPLOAD_SIZE_MB:
            raise DocumentIngestionError(
                f"File size {size_mb:.1f} MB exceeds limit of "
                f"{MAX_UPLOAD_SIZE_MB} MB (PG_MAX_UPLOAD_SIZE_MB)."
            )

        # --- deterministic doc_id from file content -----------------------
        file_bytes = file_path.read_bytes()
        doc_id = hashlib.sha256(file_bytes).hexdigest()[:16]
        logger.info(
            "[document_ingester] Ingesting %s (%.2f MB) -> doc_id=%s",
            file_path.name, size_mb, doc_id,
        )

        # --- dispatch to the appropriate parser ---------------------------
        parser_name = self.PARSERS.get(ext)
        if parser_name is None:
            raise DocumentIngestionError(
                f"No parser registered for extension '{ext}'."
            )

        parser_fn = getattr(self, parser_name, None)
        if parser_fn is None:
            raise DocumentIngestionError(
                f"Parser method '{parser_name}' not implemented."
            )

        start_ts = time.monotonic()
        result = parser_fn(file_path, file_bytes)
        elapsed = time.monotonic() - start_ts

        content: str = result["content"]
        html: str = result["html"]
        metadata: dict = result["metadata"]
        pages: int = result["pages"]

        if not content.strip():
            logger.warning(
                "[document_ingester] Extraction produced empty content "
                "for %s (doc_id=%s).",
                file_path.name, doc_id,
            )

        # --- enrich metadata ----------------------------------------------
        metadata.update({
            "doc_id": doc_id,
            "original_filename": file_path.name,
            "extension": ext,
            "size_bytes": len(file_bytes),
            "size_mb": round(size_mb, 3),
            "content_chars": len(content),
            "pages": pages,
            "parse_seconds": round(elapsed, 3),
            "ingested_at": time.time(),
        })

        # --- persist artefacts --------------------------------------------
        self._persist(doc_id, file_path, content, html, metadata)

        logger.info(
            "[document_ingester] Ingested doc_id=%s | %d chars | %d pages | %.2fs",
            doc_id, len(content), pages, elapsed,
        )

        return {
            "content": content,
            "html": html,
            "metadata": metadata,
            "pages": pages,
            "doc_id": doc_id,
        }

    # ------------------------------------------------------------------
    # Parsers (all local, no external API calls)
    # ------------------------------------------------------------------

    def _parse_pdf(self, file_path: Path, file_bytes: bytes) -> dict:
        """Extract text and HTML from a PDF using PyMuPDF (fitz).

        If the normal text extraction yields fewer than OCR_MIN_TEXT_CHARS
        non-whitespace characters (common with CorelDRAW / Illustrator PDFs
        where text is rendered as vector curves), falls back to OCR via
        pytesseract.  OCR is limited to the first OCR_MAX_PAGES pages to
        avoid timeouts on large documents.
        """
        try:
            import fitz  # PyMuPDF
        except ImportError as exc:
            raise DocumentIngestionError(
                "PyMuPDF (fitz) is required for PDF parsing. "
                "Install with: pip install PyMuPDF"
            ) from exc

        doc = fitz.open(stream=file_bytes, filetype="pdf")
        text_parts: list[str] = []
        html_parts: list[str] = ["<html><body>"]

        # Capture page count BEFORE closing the document
        page_count = len(doc)

        for page_num in range(page_count):
            page = doc[page_num]
            page_text = page.get_text("text")
            page_html = page.get_text("html")
            text_parts.append(page_text)
            html_parts.append(
                f'<div class="page" data-page="{page_num + 1}">'
                f"{page_html}</div>"
            )

        html_parts.append("</body></html>")

        pdf_metadata = doc.metadata or {}

        # -----------------------------------------------------------------
        # OCR fallback: detect vector-curve PDFs with no searchable text
        # -----------------------------------------------------------------
        combined_text = "\n\n".join(text_parts)
        non_ws_char_count = len(combined_text.replace(" ", "").replace("\n", "").replace("\r", "").replace("\t", ""))
        ocr_fallback_used = False
        ocr_pages_processed = 0

        if non_ws_char_count < OCR_MIN_TEXT_CHARS:
            logger.warning(
                "[document_ingester] PDF text extraction yielded only %d "
                "non-whitespace chars (threshold=%d) for %s. "
                "Attempting OCR fallback via pytesseract.",
                non_ws_char_count, OCR_MIN_TEXT_CHARS, file_path.name,
            )

            ocr_text_parts: list[str] = []
            ocr_html_parts: list[str] = []

            try:
                import pytesseract
                from PIL import Image

                pytesseract.pytesseract.tesseract_cmd = TESSERACT_CMD

                pages_to_ocr = min(page_count, OCR_MAX_PAGES)

                for page_num in range(pages_to_ocr):
                    page = doc[page_num]
                    pixmap = page.get_pixmap(dpi=OCR_DPI)

                    # Convert PyMuPDF pixmap to PIL Image via PNG bytes
                    png_bytes = pixmap.tobytes("png")
                    pil_image = Image.open(io.BytesIO(png_bytes))

                    ocr_page_text = pytesseract.image_to_string(pil_image)
                    ocr_text_parts.append(ocr_page_text)
                    ocr_html_parts.append(
                        f'<div class="page ocr" data-page="{page_num + 1}">'
                        f"<p>{_escape_html(ocr_page_text)}</p></div>"
                    )
                    ocr_pages_processed += 1

                ocr_combined = "\n\n".join(ocr_text_parts)
                ocr_non_ws = len(ocr_combined.replace(" ", "").replace("\n", "").replace("\r", "").replace("\t", ""))

                if ocr_non_ws > non_ws_char_count:
                    # OCR produced more text -- use it (merged with any
                    # partial text from the normal extraction)
                    combined_text = ocr_combined
                    html_parts = (
                        ["<html><body>"]
                        + ocr_html_parts
                        + ["</body></html>"]
                    )
                    ocr_fallback_used = True

                    logger.info(
                        "[document_ingester] OCR fallback succeeded for %s: "
                        "%d chars from %d pages (was %d chars from normal extraction).",
                        file_path.name, ocr_non_ws,
                        ocr_pages_processed, non_ws_char_count,
                    )
                else:
                    logger.info(
                        "[document_ingester] OCR fallback did not improve "
                        "extraction for %s (OCR=%d chars, normal=%d chars). "
                        "Keeping original extraction.",
                        file_path.name, ocr_non_ws, non_ws_char_count,
                    )

            except ImportError:
                logger.warning(
                    "[document_ingester] pytesseract or Pillow not available "
                    "for OCR fallback on %s. Install with: "
                    "pip install pytesseract Pillow",
                    file_path.name,
                )
            except Exception as exc:
                logger.warning(
                    "[document_ingester] OCR fallback failed for %s: %s. "
                    "Falling back to normal text extraction (%d chars).",
                    file_path.name, str(exc)[:300], non_ws_char_count,
                )

        doc.close()

        result_metadata = {
            "pdf_title": pdf_metadata.get("title", ""),
            "pdf_author": pdf_metadata.get("author", ""),
            "pdf_subject": pdf_metadata.get("subject", ""),
            "pdf_creator": pdf_metadata.get("creator", ""),
            "pdf_producer": pdf_metadata.get("producer", ""),
            "pdf_creation_date": pdf_metadata.get("creationDate", ""),
            "ocr_fallback_used": ocr_fallback_used,
        }

        if ocr_fallback_used:
            result_metadata["ocr_pages_processed"] = ocr_pages_processed
            result_metadata["ocr_dpi"] = OCR_DPI

        return {
            "content": combined_text,
            "html": "\n".join(html_parts),
            "metadata": result_metadata,
            "pages": page_count,
        }

    def _parse_docx(self, file_path: Path, file_bytes: bytes) -> dict:
        """Extract text from a DOCX using python-docx."""
        try:
            import docx
        except ImportError as exc:
            raise DocumentIngestionError(
                "python-docx is required for DOCX parsing. "
                "Install with: pip install python-docx"
            ) from exc

        document = docx.Document(io.BytesIO(file_bytes))

        text_parts: list[str] = []
        html_parts: list[str] = ["<html><body>"]

        # --- Core properties (metadata) ---
        core = document.core_properties
        meta = {
            "docx_title": core.title or "",
            "docx_author": core.author or "",
            "docx_subject": core.subject or "",
            "docx_created": str(core.created) if core.created else "",
            "docx_modified": str(core.modified) if core.modified else "",
        }

        # --- Paragraphs ---
        for para in document.paragraphs:
            text = para.text.strip()
            if not text:
                continue
            text_parts.append(text)

            style_name = para.style.name if para.style else ""
            if style_name.startswith("Heading"):
                # Attempt to determine heading level
                level = style_name.replace("Heading", "").strip()
                level = level if level.isdigit() else "2"
                html_parts.append(f"<h{level}>{_escape_html(text)}</h{level}>")
            else:
                html_parts.append(f"<p>{_escape_html(text)}</p>")

        # --- Tables ---
        for table_idx, table in enumerate(document.tables):
            table_text_rows: list[str] = []
            html_parts.append("<table border='1'>")
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                table_text_rows.append(" | ".join(cells))
                html_parts.append(
                    "<tr>"
                    + "".join(f"<td>{_escape_html(c)}</td>" for c in cells)
                    + "</tr>"
                )
            html_parts.append("</table>")
            if table_text_rows:
                text_parts.append(
                    f"\n[Table {table_idx + 1}]\n"
                    + "\n".join(table_text_rows)
                )

        html_parts.append("</body></html>")

        # Estimate "pages" from paragraph count (DOCX has no native pages).
        estimated_pages = max(1, len(text_parts) // 40)

        return {
            "content": "\n\n".join(text_parts),
            "html": "\n".join(html_parts),
            "metadata": meta,
            "pages": estimated_pages,
        }

    def _parse_doc(self, file_path: Path, file_bytes: bytes) -> dict:
        """Extract text from a legacy .doc file.

        Strategy:
        1. Try antiword (best quality for binary .doc files).
        2. Fall back to raw text extraction with error suppression.

        Legacy .doc (OLE2/CFBF) is a binary format; without a dedicated
        parser the best we can do is extract printable text sequences.
        """
        logger.info(
            "[document_ingester] Parsing legacy .doc file: %s",
            file_path.name,
        )

        content = ""
        extraction_method = "none"

        # --- Strategy 1: antiword -----------------------------------------
        try:
            result = subprocess.run(
                [ANTIWORD_CMD, str(file_path)],
                capture_output=True,
                text=True,
                timeout=30,
            )
            if result.returncode == 0 and result.stdout.strip():
                content = result.stdout.strip()
                extraction_method = "antiword"
                logger.info(
                    "[document_ingester] antiword extracted %d chars from %s",
                    len(content), file_path.name,
                )
        except FileNotFoundError:
            logger.info(
                "[document_ingester] antiword not available on this system. "
                "Falling back to raw text extraction for %s.",
                file_path.name,
            )
        except subprocess.TimeoutExpired:
            logger.warning(
                "[document_ingester] antiword timed out on %s.",
                file_path.name,
            )
        except Exception as exc:
            logger.warning(
                "[document_ingester] antiword failed for %s: %s",
                file_path.name, str(exc)[:200],
            )

        # --- Strategy 2: raw text extraction (fallback) -------------------
        if not content:
            logger.info(
                "[document_ingester] Attempting raw text extraction for %s",
                file_path.name,
            )
            try:
                # Try UTF-8, then latin-1 for binary .doc files
                try:
                    raw_text = file_bytes.decode("utf-8", errors="ignore")
                except Exception:
                    raw_text = file_bytes.decode("latin-1", errors="ignore")

                # Filter to printable text sequences (at least 4 chars long)
                import re
                # Extract runs of printable ASCII/Latin characters
                text_runs = re.findall(r"[\x20-\x7E\xA0-\xFF]{4,}", raw_text)
                content = "\n".join(text_runs)
                extraction_method = "raw_text"

                logger.info(
                    "[document_ingester] Raw text extraction yielded %d chars "
                    "from %d text runs for %s",
                    len(content), len(text_runs), file_path.name,
                )
            except Exception as exc:
                logger.warning(
                    "[document_ingester] Raw text extraction failed for %s: %s",
                    file_path.name, str(exc)[:200],
                )
                content = f"[Binary .doc file: {file_path.name} -- text extraction unavailable]"
                extraction_method = "failed"

        html = _text_to_html(content, is_markdown=False)

        return {
            "content": content,
            "html": html,
            "metadata": {
                "doc_format": "legacy_doc",
                "extraction_method": extraction_method,
            },
            "pages": 1,
        }

    def _parse_xlsx(self, file_path: Path, file_bytes: bytes) -> dict:
        """Extract cell values from an XLSX using openpyxl."""
        try:
            import openpyxl
        except ImportError as exc:
            raise DocumentIngestionError(
                "openpyxl is required for XLSX parsing. "
                "Install with: pip install openpyxl"
            ) from exc

        workbook = openpyxl.load_workbook(
            io.BytesIO(file_bytes),
            data_only=True,
            read_only=True,
        )

        text_parts: list[str] = []
        html_parts: list[str] = ["<html><body>"]
        sheet_count = 0

        # Capture sheet names BEFORE closing (read_only workbooks may
        # not support attribute access after close).
        sheet_names = list(workbook.sheetnames)

        for sheet_name in sheet_names:
            sheet = workbook[sheet_name]
            sheet_count += 1

            text_parts.append(f"[Sheet: {sheet_name}]")
            html_parts.append(f"<h2>{_escape_html(sheet_name)}</h2>")
            html_parts.append("<table border='1'>")

            for row in sheet.iter_rows():
                cell_values = []
                for cell in row:
                    val = cell.value
                    cell_values.append(str(val) if val is not None else "")

                # Skip entirely empty rows
                if not any(v.strip() for v in cell_values):
                    continue

                text_parts.append(" | ".join(cell_values))
                html_parts.append(
                    "<tr>"
                    + "".join(
                        f"<td>{_escape_html(v)}</td>" for v in cell_values
                    )
                    + "</tr>"
                )

            html_parts.append("</table>")

        html_parts.append("</body></html>")
        workbook.close()

        return {
            "content": "\n".join(text_parts),
            "html": "\n".join(html_parts),
            "metadata": {
                "xlsx_sheet_count": sheet_count,
                "xlsx_sheet_names": sheet_names,
            },
            "pages": sheet_count,
        }

    def _parse_pptx(self, file_path: Path, file_bytes: bytes) -> dict:
        """Extract slide text from a PPTX using python-pptx."""
        try:
            from pptx import Presentation
        except ImportError as exc:
            raise DocumentIngestionError(
                "python-pptx is required for PPTX parsing. "
                "Install with: pip install python-pptx"
            ) from exc

        prs = Presentation(io.BytesIO(file_bytes))

        text_parts: list[str] = []
        html_parts: list[str] = ["<html><body>"]
        slide_count = 0

        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_count += 1
            slide_texts: list[str] = []

            html_parts.append(
                f'<div class="slide" data-slide="{slide_num}">'
                f"<h3>Slide {slide_num}</h3>"
            )

            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    para_text = paragraph.text.strip()
                    if para_text:
                        slide_texts.append(para_text)
                        html_parts.append(f"<p>{_escape_html(para_text)}</p>")

            # Include slide notes if present
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    slide_texts.append(f"[Notes: {notes_text}]")
                    html_parts.append(
                        f"<p><em>[Notes: {_escape_html(notes_text)}]</em></p>"
                    )

            html_parts.append("</div>")

            if slide_texts:
                text_parts.append(
                    f"[Slide {slide_num}]\n" + "\n".join(slide_texts)
                )

        html_parts.append("</body></html>")

        return {
            "content": "\n\n".join(text_parts),
            "html": "\n".join(html_parts),
            "metadata": {
                "pptx_slide_count": slide_count,
            },
            "pages": slide_count,
        }

    def _parse_text(self, file_path: Path, file_bytes: bytes) -> dict:
        """Read plain text or Markdown files directly."""
        # Try UTF-8 first, fall back to latin-1
        try:
            content = file_bytes.decode("utf-8")
        except UnicodeDecodeError:
            content = file_bytes.decode("latin-1")
            logger.info(
                "[document_ingester] Fell back to latin-1 encoding for %s",
                file_path.name,
            )

        ext = file_path.suffix.lower()
        html = _text_to_html(content, is_markdown=(ext == ".md"))

        return {
            "content": content,
            "html": html,
            "metadata": {
                "encoding": "utf-8",
                "line_count": content.count("\n") + 1,
            },
            "pages": 1,
        }

    def _parse_csv(self, file_path: Path, file_bytes: bytes) -> dict:
        """Parse CSV into a text table using the stdlib csv module."""
        try:
            text = file_bytes.decode("utf-8")
        except UnicodeDecodeError:
            text = file_bytes.decode("latin-1")

        reader = csv.reader(io.StringIO(text))
        rows: list[list[str]] = []
        for row in reader:
            rows.append(row)

        if not rows:
            return {
                "content": "",
                "html": "<html><body><p>Empty CSV</p></body></html>",
                "metadata": {"csv_rows": 0, "csv_columns": 0},
                "pages": 1,
            }

        # Determine column widths for alignment
        col_count = max(len(r) for r in rows)
        col_widths = [0] * col_count
        for row in rows:
            for i, cell in enumerate(row):
                if i < col_count:
                    col_widths[i] = max(col_widths[i], len(cell))

        # Build plain text table
        text_lines: list[str] = []
        for row_idx, row in enumerate(rows):
            padded = []
            for i in range(col_count):
                val = row[i] if i < len(row) else ""
                padded.append(val.ljust(col_widths[i]))
            text_lines.append(" | ".join(padded))
            # Separator after header row
            if row_idx == 0:
                text_lines.append(
                    " | ".join("-" * w for w in col_widths)
                )

        # Build HTML table
        html_parts: list[str] = ["<html><body><table border='1'>"]
        for row_idx, row in enumerate(rows):
            tag = "th" if row_idx == 0 else "td"
            cells_html = "".join(
                f"<{tag}>{_escape_html(row[i] if i < len(row) else '')}</{tag}>"
                for i in range(col_count)
            )
            html_parts.append(f"<tr>{cells_html}</tr>")
        html_parts.append("</table></body></html>")

        return {
            "content": "\n".join(text_lines),
            "html": "\n".join(html_parts),
            "metadata": {
                "csv_rows": len(rows),
                "csv_columns": col_count,
            },
            "pages": 1,
        }

    def _parse_html(self, file_path: Path, file_bytes: bytes) -> dict:
        """Extract article content from HTML using readability-lxml."""
        try:
            text = file_bytes.decode("utf-8")
        except UnicodeDecodeError:
            text = file_bytes.decode("latin-1")

        title = ""
        article_html = text
        article_text = ""

        try:
            from readability import Document as ReadabilityDocument

            doc = ReadabilityDocument(text)
            title = doc.title() or ""
            article_html = doc.summary() or text
        except ImportError:
            logger.warning(
                "[document_ingester] readability-lxml not installed. "
                "Falling back to raw HTML. "
                "Install with: pip install readability-lxml"
            )
        except Exception as exc:
            logger.warning(
                "[document_ingester] readability extraction failed for %s: %s. "
                "Falling back to raw HTML.",
                file_path.name, str(exc)[:200],
            )

        # Extract plain text from the (possibly cleaned) HTML
        try:
            from lxml import etree

            tree = etree.HTML(article_html)
            if tree is not None:
                article_text = " ".join(tree.itertext()).strip()
        except ImportError:
            logger.warning(
                "[document_ingester] lxml not installed for text extraction. "
                "Using basic tag stripping."
            )
            import re
            article_text = re.sub(r"<[^>]+>", " ", article_html)
            article_text = re.sub(r"\s+", " ", article_text).strip()
        except Exception as exc:
            logger.warning(
                "[document_ingester] lxml text extraction failed: %s",
                str(exc)[:200],
            )
            import re
            article_text = re.sub(r"<[^>]+>", " ", article_html)
            article_text = re.sub(r"\s+", " ", article_text).strip()

        return {
            "content": article_text,
            "html": article_html,
            "metadata": {
                "html_title": title,
            },
            "pages": 1,
        }

    def _parse_image(self, file_path: Path, file_bytes: bytes) -> dict:
        """Extract text from an image file via OCR (pytesseract).

        Also captures image dimensions and format via PIL for metadata.
        If OCR yields no text, returns a placeholder description.
        """
        logger.info(
            "[document_ingester] Parsing image file via OCR: %s",
            file_path.name,
        )

        content = ""
        width = 0
        height = 0
        image_format = file_path.suffix.lower().lstrip(".")

        try:
            from PIL import Image

            pil_image = Image.open(io.BytesIO(file_bytes))
            width = pil_image.width
            height = pil_image.height
            image_format = pil_image.format or image_format

            logger.info(
                "[document_ingester] Image %s: %dx%d, format=%s",
                file_path.name, width, height, image_format,
            )

            # Attempt OCR via pytesseract
            try:
                import pytesseract

                pytesseract.pytesseract.tesseract_cmd = TESSERACT_CMD

                # Convert to RGB if necessary (e.g., RGBA PNGs, palette GIFs)
                if pil_image.mode not in ("L", "RGB"):
                    pil_image = pil_image.convert("RGB")

                ocr_text = pytesseract.image_to_string(pil_image)
                if ocr_text and ocr_text.strip():
                    content = ocr_text.strip()
                    logger.info(
                        "[document_ingester] OCR extracted %d chars from %s",
                        len(content), file_path.name,
                    )
                else:
                    logger.info(
                        "[document_ingester] OCR yielded no text for %s",
                        file_path.name,
                    )
            except ImportError:
                logger.warning(
                    "[document_ingester] pytesseract not available for OCR on %s. "
                    "Install with: pip install pytesseract",
                    file_path.name,
                )
            except Exception as exc:
                logger.warning(
                    "[document_ingester] OCR failed for %s: %s",
                    file_path.name, str(exc)[:300],
                )

        except ImportError:
            logger.warning(
                "[document_ingester] Pillow not available for image parsing. "
                "Install with: pip install Pillow"
            )
        except Exception as exc:
            logger.warning(
                "[document_ingester] Failed to open image %s: %s",
                file_path.name, str(exc)[:300],
            )

        # Fallback placeholder if OCR produced nothing
        if not content:
            content = f"[Image: {file_path.name}]"

        html = (
            f"<html><body>"
            f'<div class="image-content">'
            f"<p><strong>Image:</strong> {_escape_html(file_path.name)} "
            f"({width}x{height} {_escape_html(image_format)})</p>"
            f"<pre>{_escape_html(content)}</pre>"
            f"</div>"
            f"</body></html>"
        )

        return {
            "content": content,
            "html": html,
            "metadata": {
                "image_width": width,
                "image_height": height,
                "image_format": image_format,
            },
            "pages": 1,
        }

    def _parse_audio(self, file_path: Path, file_bytes: bytes) -> dict:
        """Transcribe an audio file using OpenAI Whisper (local model).

        Uses the Whisper model specified by WHISPER_MODEL env var (default: "base").
        Falls back gracefully with a warning if whisper is unavailable.
        """
        logger.info(
            "[document_ingester] Parsing audio file via Whisper: %s",
            file_path.name,
        )

        audio_format = file_path.suffix.lower().lstrip(".")
        content = ""
        duration_seconds = 0.0
        detected_language = ""

        try:
            import whisper

            # Set ffmpeg path if configured (Whisper uses ffmpeg internally)
            if FFMPEG_PATH and os.path.isfile(FFMPEG_PATH):
                os.environ.setdefault("PATH", "")
                ffmpeg_dir = str(Path(FFMPEG_PATH).parent)
                if ffmpeg_dir not in os.environ["PATH"]:
                    os.environ["PATH"] = ffmpeg_dir + os.pathsep + os.environ["PATH"]

            model_name = WHISPER_MODEL
            logger.info(
                "[document_ingester] Loading Whisper model '%s' for %s",
                model_name, file_path.name,
            )
            model = whisper.load_model(model_name)

            # Whisper needs a file path on disk (it uses ffmpeg to decode).
            # file_path already points to the uploaded file on disk.
            result = model.transcribe(str(file_path))
            content = result.get("text", "").strip()
            detected_language = result.get("language", "")

            # Extract duration from segments if available
            segments = result.get("segments", [])
            if segments:
                last_segment = segments[-1]
                duration_seconds = last_segment.get("end", 0.0)

            logger.info(
                "[document_ingester] Whisper transcribed %d chars, "
                "language=%s, duration=%.1fs from %s",
                len(content), detected_language,
                duration_seconds, file_path.name,
            )

        except ImportError:
            logger.warning(
                "[document_ingester] openai-whisper not available for audio "
                "transcription on %s. "
                "Install with: pip install openai-whisper",
                file_path.name,
            )
            content = f"[Audio: {file_path.name} -- whisper not installed for transcription]"
        except Exception as exc:
            logger.warning(
                "[document_ingester] Whisper transcription failed for %s: %s",
                file_path.name, str(exc)[:300],
            )
            content = f"[Audio: {file_path.name} -- transcription failed]"

        html = (
            f"<html><body>"
            f'<div class="audio-content">'
            f"<p><strong>Audio Transcription:</strong> {_escape_html(file_path.name)} "
            f"({_escape_html(audio_format)}"
        )
        if duration_seconds > 0:
            minutes = int(duration_seconds // 60)
            seconds = int(duration_seconds % 60)
            html += f", {minutes}m {seconds}s"
        if detected_language:
            html += f", {_escape_html(detected_language)}"
        html += (
            f")</p>"
            f"<pre>{_escape_html(content)}</pre>"
            f"</div>"
            f"</body></html>"
        )

        return {
            "content": content,
            "html": html,
            "metadata": {
                "audio_format": audio_format,
                "duration_seconds": round(duration_seconds, 2),
                "language": detected_language,
            },
            "pages": 1,
        }

    # ------------------------------------------------------------------
    # Persistence
    # ------------------------------------------------------------------

    def _persist(
        self,
        doc_id: str,
        original_path: Path,
        content: str,
        html: str,
        metadata: dict,
    ) -> None:
        """Store extracted artefacts to disk under data/documents/{doc_id}/."""
        doc_dir = DOCUMENT_STORAGE_DIR / doc_id
        doc_dir.mkdir(parents=True, exist_ok=True)

        # Copy original file
        dest_original = doc_dir / f"original{original_path.suffix.lower()}"
        shutil.copy2(str(original_path), str(dest_original))

        # Write extracted text
        extracted_txt = doc_dir / "extracted.txt"
        extracted_txt.write_text(content, encoding="utf-8")

        # Write extracted HTML
        extracted_html = doc_dir / "extracted.html"
        extracted_html.write_text(html, encoding="utf-8")

        # Write metadata JSON
        metadata_json = doc_dir / "metadata.json"
        metadata_json.write_text(
            json.dumps(metadata, indent=2, ensure_ascii=False, default=str),
            encoding="utf-8",
        )

        logger.info(
            "[document_ingester] Persisted artefacts to %s "
            "(txt=%d bytes, html=%d bytes)",
            doc_dir, len(content), len(html),
        )

    # ------------------------------------------------------------------
    # Utility: list ingested documents
    # ------------------------------------------------------------------

    def list_documents(self) -> list[dict]:
        """Return metadata for all previously ingested documents."""
        results: list[dict] = []
        if not DOCUMENT_STORAGE_DIR.exists():
            return results

        for child in sorted(DOCUMENT_STORAGE_DIR.iterdir()):
            if not child.is_dir():
                continue
            meta_file = child / "metadata.json"
            if meta_file.exists():
                try:
                    meta = json.loads(meta_file.read_text(encoding="utf-8"))
                    results.append(meta)
                except Exception as exc:
                    logger.warning(
                        "[document_ingester] Failed to read metadata "
                        "for %s: %s",
                        child.name, str(exc)[:200],
                    )
        return results

    def get_document(self, doc_id: str) -> Optional[dict]:
        """Load a previously ingested document by its doc_id.

        Returns dict with content, html, metadata, pages, doc_id
        or None if not found.
        """
        doc_dir = DOCUMENT_STORAGE_DIR / doc_id
        if not doc_dir.exists():
            return None

        meta_path = doc_dir / "metadata.json"
        txt_path = doc_dir / "extracted.txt"
        html_path = doc_dir / "extracted.html"

        if not meta_path.exists():
            logger.warning(
                "[document_ingester] Missing metadata.json for doc_id=%s",
                doc_id,
            )
            return None

        metadata = json.loads(meta_path.read_text(encoding="utf-8"))
        content = txt_path.read_text(encoding="utf-8") if txt_path.exists() else ""
        html = html_path.read_text(encoding="utf-8") if html_path.exists() else ""

        return {
            "content": content,
            "html": html,
            "metadata": metadata,
            "pages": metadata.get("pages", 1),
            "doc_id": doc_id,
        }


# ---------------------------------------------------------------------------
# Module-level helpers
# ---------------------------------------------------------------------------

def _escape_html(text: str) -> str:
    """Escape HTML special characters."""
    return (
        text
        .replace("&", "&amp;")
        .replace("<", "&lt;")
        .replace(">", "&gt;")
        .replace('"', "&quot;")
        .replace("'", "&#x27;")
    )


def _text_to_html(content: str, is_markdown: bool = False) -> str:
    """Convert plain text or Markdown to simple HTML.

    For Markdown, headings (# lines) and code fences are converted.
    For plain text, each line becomes a <p> tag.
    """
    lines = content.split("\n")
    html_parts: list[str] = ["<html><body>"]

    in_code_block = False

    for line in lines:
        stripped = line.strip()

        # Code fences (Markdown)
        if is_markdown and stripped.startswith("```"):
            if in_code_block:
                html_parts.append("</pre></code>")
                in_code_block = False
            else:
                html_parts.append("<code><pre>")
                in_code_block = True
            continue

        if in_code_block:
            html_parts.append(_escape_html(line))
            continue

        if not stripped:
            continue

        # Markdown headings
        if is_markdown and stripped.startswith("#"):
            level = 0
            for ch in stripped:
                if ch == "#":
                    level += 1
                else:
                    break
            level = min(level, 6)
            heading_text = stripped[level:].strip()
            html_parts.append(
                f"<h{level}>{_escape_html(heading_text)}</h{level}>"
            )
        else:
            html_parts.append(f"<p>{_escape_html(stripped)}</p>")

    if in_code_block:
        html_parts.append("</pre></code>")

    html_parts.append("</body></html>")
    return "\n".join(html_parts)
